"""
Generate a fully editable PowerPoint deck for the SQL Agent class
presentation. Uses python-pptx so every text element is a native shape
that can be clicked into and edited in PowerPoint or Keynote.

Run:
    python3 generate_pptx.py
Output: slides-editable.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------------------------------------------------------------- theme
INK = RGBColor(0x0E, 0x0E, 0x0E)
INK_MUTED = RGBColor(0x5A, 0x5A, 0x5A)
INK_FAINT = RGBColor(0xE5, 0xE5, 0xE5)
SURFACE = RGBColor(0xFA, 0xFA, 0xF9)
SURFACE_RAISED = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC9, 0x64, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"  # Closest cross-platform readable to SF Pro
FONT_MONO = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ----------------------------------------------------------- helpers
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, x, y, w, h, paragraphs, size=14, color=INK, font=FONT,
                   line_spacing=1.45):
    """paragraphs: list of either strings or list of (text, opts) tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(4)
        if isinstance(para, str):
            run = p.add_run()
            run.text = para
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
        else:
            # list of runs: each is (text, {bold, color, font, size})
            for text, opts in para:
                run = p.add_run()
                run.text = text
                run.font.name = opts.get("font", font)
                run.font.size = Pt(opts.get("size", size))
                run.font.bold = opts.get("bold", False)
                run.font.color.rgb = opts.get("color", color)
    return box


def add_rect(slide, x, y, w, h, fill=SURFACE_RAISED, line=INK_FAINT,
             line_w=0.75, corner=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color=INK_FAINT, weight=0.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_accent_bar(slide):
    """Small amber bar at top-left, the project's accent mark."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.83), Inches(0.45),
        Inches(0.32), Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_kicker(slide, text, x=Inches(0.83), y=Inches(0.6)):
    add_text(slide, x, y, Inches(8), Inches(0.3),
             text.upper(), size=10, bold=True, color=ACCENT,
             line_spacing=1.0)


def add_title(slide, text, x=Inches(0.83), y=Inches(0.95), w=Inches(11.67),
              h=Inches(1.0), size=28):
    add_text(slide, x, y, w, h, text, size=size, bold=True, color=INK,
             line_spacing=1.15)


def add_footer(slide, label, page_num):
    add_text(slide, Inches(0.83), Inches(7.05), Inches(8), Inches(0.3),
             label, size=9, color=INK_MUTED)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.5), Inches(0.3),
             str(page_num), size=9, color=INK_MUTED, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------- build
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]


# ============================================================ 01 TITLE
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)

add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.5),
         "SQL Agent", size=64, bold=True, color=INK, line_spacing=1.05)
add_text(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.6),
         "A multi-model system for natural-language data analysis",
         size=20, color=INK_MUTED)

add_line(slide, Inches(1.5), Inches(5.3), Inches(7), Inches(5.3),
         color=INK_FAINT, weight=0.5)

add_text(slide, Inches(1.5), Inches(5.45), Inches(2), Inches(0.3),
         "TEAM", size=10, bold=True, color=INK_MUTED)
add_text(slide, Inches(1.5), Inches(5.75), Inches(10), Inches(0.4),
         "Daniel Regalado Cardoso · Nefeli Zafeiri · Oliver Mazariegos · Eleniz Espina",
         size=14, color=INK)
add_text(slide, Inches(1.5), Inches(6.1), Inches(10), Inches(0.4),
         "MSBA · University of Miami · 2026",
         size=12, color=INK_MUTED)


# ============================================================ 02 PROBLEM
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Problem statement")
add_title(slide, "Tabular data analysis still requires SQL knowledge")

# Left column
add_text(slide, Inches(0.83), Inches(2.4), Inches(5.5), Inches(0.4),
         "Business users frequently have:", size=15, color=INK)
add_paragraphs(slide, Inches(0.83), Inches(2.85), Inches(5.5), Inches(2.5), [
    "•  Tabular data they need to query",
    "•  Specific questions in mind",
    "•  No working SQL knowledge",
], size=14, color=INK)
add_text(slide, Inches(0.83), Inches(4.6), Inches(5.5), Inches(1.5),
         "Existing options (manual SQL, generic chatbots, BI tools) each have trade-offs in cost, accuracy, or accessibility.",
         size=13, color=INK_MUTED)

# Right column — Objective card
card = add_rect(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.0))
add_text(slide, Inches(7.3), Inches(2.55), Inches(5), Inches(0.3),
         "OBJECTIVE", size=10, bold=True, color=ACCENT)
add_text(slide, Inches(7.3), Inches(2.85), Inches(5), Inches(0.5),
         "Build a system that:", size=15, bold=True, color=INK)
add_paragraphs(slide, Inches(7.3), Inches(3.4), Inches(5), Inches(3), [
    "1.  Accepts a CSV or JSON file as input",
    "2.  Accepts a question in natural language",
    "3.  Returns the answer as a chart and a written finding",
    "4.  Runs at low cost on free GPU infrastructure",
], size=13, color=INK)

add_footer(slide, "Problem statement", "02 / 13")


# ============================================================ 03 METHODOLOGY
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Methodology")
add_title(slide, "Six phases from raw data to deployed system")

# Six boxes left to right
phases = [
    ("1", "Source", "10 public\ndatasets"),
    ("2", "Curate", "1.2M to 723k\nrows"),
    ("3", "Build", "3 task\ndatasets"),
    ("4", "Train", "3 LoRA\nadapters"),
    ("5", "Architect", "Orchestrator\n+ DuckDB"),
    ("6", "Deploy", "Hugging Face\nSpaces"),
]
box_w = Inches(1.85)
gap = Inches(0.15)
total_w = box_w * 6 + gap * 5
start_x = (SLIDE_W - total_w) / 2

for i, (num, name, sub) in enumerate(phases):
    x = start_x + (box_w + gap) * i
    accent = i >= 3
    line_color = ACCENT if accent else INK_FAINT
    line_w = 1.5 if accent else 0.75
    rect = add_rect(slide, x, Inches(3.2), box_w, Inches(1.6),
                    fill=SURFACE_RAISED, line=line_color, line_w=line_w)
    add_text(slide, x, Inches(3.35), box_w, Inches(0.3), num,
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.6), box_w, Inches(0.4), name,
             size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.0), box_w, Inches(0.7), sub,
             size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)
    # arrow
    if i < 5:
        ax = x + box_w + Inches(0.02)
        add_line(slide, ax, Inches(4.0), ax + Inches(0.11), Inches(4.0),
                 color=INK_MUTED, weight=0.75)

add_text(slide, Inches(0.83), Inches(5.5), Inches(11.67), Inches(0.6),
         "Each phase produces a reproducible artifact — a dataset, a model adapter, or a deployable component.",
         size=13, color=INK_MUTED)
add_text(slide, Inches(0.83), Inches(5.95), Inches(11.67), Inches(0.6),
         "The repository contains the scripts to recreate each step.",
         size=13, color=INK_MUTED)

add_footer(slide, "Methodology overview", "03 / 13")


# ============================================================ 04 PHASE 1 SOURCING
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 1 — Sourcing")
add_title(slide, "Aggregating ten public text-to-SQL datasets")

# Left list
add_text(slide, Inches(0.83), Inches(2.4), Inches(5.5), Inches(0.4),
         "Selected ten existing datasets covering different schemas and SQL dialects:",
         size=13, color=INK)
sources = [
    "b-mc2/sql-create-context", "gretelai/synthetic_text_to_sql",
    "knowrohit07/know_sql", "NumbersStation/NSText2SQL",
    "Clinton/Text-to-sql-v1", "motherduckdb/duckdb-text2sql-25k",
    "bugdaryan/spider-natsql-wikisql", "ChrisHayduk/Llama-2-SQL",
    "kaxap/llama2-sql-instruct", "PipableAI/spider-bird",
]
add_paragraphs(slide, Inches(0.83), Inches(3.0), Inches(5.5), Inches(4),
               [f"•  {s}" for s in sources], size=12, color=INK, font=FONT_MONO,
               line_spacing=1.5)

# Right column
add_text(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.4),
         "RATIONALE", size=10, bold=True, color=ACCENT)
add_text(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(1.6),
         "Building a corpus from scratch was not feasible within the project timeline. Public datasets give us schema diversity and large coverage with permissive licenses.",
         size=13, color=INK)
add_text(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(0.4),
         "COMBINED SIZE", size=10, bold=True, color=ACCENT)
add_text(slide, Inches(7.0), Inches(4.9), Inches(5.5), Inches(1.6),
         "Approximately 1.2 million rows before cleaning. Heterogeneous formats and varying quality, requiring substantial pre-processing.",
         size=13, color=INK)

add_footer(slide, "Phase 1 · Sourcing", "04 / 13")


# ============================================================ 05 PHASE 2 CURATION
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 2 — Curation")
add_title(slide, "Cleaning the merged corpus")

# Pipeline blocks
steps = [
    ("10 raw sources", "1.2M rows", False),
    ("Schema unification", "7-col format", False),
    ("Deduplication", "question hashing", False),
    ("Length filter", "≤ 1024 tokens", False),
    ("Final corpus", "761,155 rows", True),
    ("Train / Val / Test", "723k / 19k / 19k", True),
]
sw = Inches(1.85); sg = Inches(0.12)
total = sw * 6 + sg * 5
sx = (SLIDE_W - total) / 2
for i, (top, bot, accent) in enumerate(steps):
    x = sx + (sw + sg) * i
    line_color = ACCENT if accent else INK_FAINT
    line_w = 1.5 if accent else 0.75
    add_rect(slide, x, Inches(2.7), sw, Inches(1.4), fill=SURFACE_RAISED,
             line=line_color, line_w=line_w)
    add_text(slide, x, Inches(2.95), sw, Inches(0.5), top,
             size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.5), sw, Inches(0.5), bot,
             size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)
    if i < 5:
        ax = x + sw + Inches(0.02)
        add_line(slide, ax, Inches(3.4), ax + Inches(0.08), Inches(3.4),
                 color=INK_MUTED, weight=0.75)

# Stats below
stats = [("1.2M", "Raw rows"), ("761k", "After dedup"),
         ("93%", "Pass length filter"), ("723k", "Final training set")]
sw = Inches(2.5); sg = Inches(0.4)
total = sw * 4 + sg * 3
sx = (SLIDE_W - total) / 2
for i, (val, label) in enumerate(stats):
    x = sx + (sw + sg) * i
    add_text(slide, x, Inches(4.7), sw, Inches(0.7), val,
             size=36, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.55), sw, Inches(0.4), label.upper(),
             size=10, color=INK_MUTED, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.83), Inches(6.4), Inches(11.67), Inches(0.4),
         "All transformations are implemented as UV scripts in training/data_pipelines/ and are reproducible.",
         size=11, color=INK_MUTED)

add_footer(slide, "Phase 2 · Curation", "05 / 13")


# ============================================================ 06 PHASE 3 DATASETS
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 3 — Task datasets")
add_title(slide, "Three datasets, one per model task")

datasets = [
    ("text-to-sql-mix-v2", "761,155", "Apache-2.0",
     "10 merged public sources",
     "huggingface.co/datasets/DanielRegaladoCardoso/text-to-sql-mix-v2"),
    ("chart-reasoning-mix-v1", "~75,000", "CC-BY-4.0",
     "nvBench (25k) + GPT-4.1-nano knowledge distillation (50k)",
     "huggingface.co/datasets/DanielRegaladoCardoso/chart-reasoning-mix-v1"),
    ("svg-chart-render-v1", "~25,000", "Apache-2.0",
     "nvBench charts re-rendered via matplotlib SVG backend",
     "huggingface.co/datasets/DanielRegaladoCardoso/svg-chart-render-v1"),
]

cw = Inches(3.95); cg = Inches(0.25)
total = cw * 3 + cg * 2
cx = (SLIDE_W - total) / 2
for i, (name, rows, lic, source, url) in enumerate(datasets):
    x = cx + (cw + cg) * i
    add_rect(slide, x, Inches(2.3), cw, Inches(4.3),
             fill=SURFACE_RAISED, line=INK_FAINT, line_w=0.75)
    # accent stripe top
    accent_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.15), Inches(2.45),
        Inches(0.3), Inches(0.05))
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = ACCENT
    accent_strip.line.fill.background()

    add_text(slide, x + Inches(0.15), Inches(2.6), cw - Inches(0.3), Inches(0.3),
             f"DATASET · 0{i+1}", size=10, bold=True, color=ACCENT)
    add_text(slide, x + Inches(0.15), Inches(2.95), cw - Inches(0.3), Inches(0.5),
             name, size=18, bold=True, color=INK, font=FONT_MONO)
    add_text(slide, x + Inches(0.15), Inches(3.6), cw - Inches(0.3), Inches(0.7),
             rows, size=30, bold=True, color=INK)
    add_text(slide, x + Inches(0.15), Inches(4.3), cw - Inches(0.3), Inches(0.4),
             "ROWS", size=9, color=INK_MUTED)
    add_text(slide, x + Inches(0.15), Inches(4.7), cw - Inches(0.3), Inches(0.4),
             f"License · {lic}", size=11, color=INK_MUTED)
    add_text(slide, x + Inches(0.15), Inches(5.0), cw - Inches(0.3), Inches(1.0),
             source, size=11, color=INK)
    add_text(slide, x + Inches(0.15), Inches(6.05), cw - Inches(0.3), Inches(0.4),
             url, size=9, color=ACCENT)

add_footer(slide, "Phase 3 · Task datasets", "06 / 13")


# ============================================================ 07 PHASE 4 SETUP
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 4 — Training setup")
add_title(slide, "QLoRA via Unsloth")

# Left column
add_text(slide, Inches(0.83), Inches(2.4), Inches(6.0), Inches(0.6),
         "We use 4-bit QLoRA with the Unsloth library for all three fine-tunes.",
         size=14, color=INK)
add_text(slide, Inches(0.83), Inches(3.1), Inches(6.0), Inches(0.4),
         "REASONS", size=10, bold=True, color=ACCENT)
add_paragraphs(slide, Inches(0.83), Inches(3.4), Inches(6.0), Inches(3.5), [
    "•  Allows training a 7B model on a single 48 GB GPU",
    "•  Approximately 2× faster than vanilla transformers",
    "•  Approximately 40% lower memory consumption",
    "•  Output is a 160 MB adapter rather than 14 GB of full weights",
], size=13, color=INK)

# Right comparison table
tx = Inches(7.4); ty = Inches(2.4); tw = Inches(5.0); th = Inches(3.5)
add_rect(slide, tx, ty, tw, th, fill=SURFACE_RAISED, line=INK_FAINT, line_w=0.75)

# Table headers
add_text(slide, tx + Inches(0.2), ty + Inches(0.2), Inches(2.0), Inches(0.3),
         "ASPECT", size=9, bold=True, color=INK_MUTED)
add_text(slide, tx + Inches(2.2), ty + Inches(0.2), Inches(1.4), Inches(0.3),
         "VANILLA", size=9, bold=True, color=INK_MUTED)
add_text(slide, tx + Inches(3.6), ty + Inches(0.2), Inches(1.4), Inches(0.3),
         "UNSLOTH", size=9, bold=True, color=INK_MUTED)

rows = [
    ("7B on 48 GB GPU", "does not fit", "fits in 4-bit"),
    ("Speed", "baseline", "~2× faster"),
    ("Memory", "baseline", "~40% less"),
    ("Output", "14 GB", "160 MB adapter"),
]
for i, (a, b, c) in enumerate(rows):
    yy = ty + Inches(0.6 + i * 0.5)
    add_text(slide, tx + Inches(0.2), yy, Inches(2.0), Inches(0.4), a, size=12, color=INK)
    add_text(slide, tx + Inches(2.2), yy, Inches(1.4), Inches(0.4), b, size=12, color=INK_MUTED)
    add_text(slide, tx + Inches(3.6), yy, Inches(1.4), Inches(0.4), c, size=12, bold=True, color=INK)

add_footer(slide, "Phase 4 · Training setup", "07 / 13")


# ============================================================ 08-10 THREE MODELS
def model_slide(num_label, page, role_num, name, function, base, dataset, dataset_url,
                method, hardware, time, loss, cost, size, hub, note):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, SURFACE)
    add_accent_bar(slide)
    add_kicker(slide, f"Phase 4 — Three fine-tunes ({role_num} of 3)")
    add_title(slide, f"Model 0{role_num} · {name}")

    # Left card: function + setup
    lx = Inches(0.83); ly = Inches(2.3); lw = Inches(6.0); lh = Inches(4.5)
    add_rect(slide, lx, ly, lw, lh, fill=SURFACE_RAISED, line=INK_FAINT)

    # accent stripe top
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                lx + Inches(0.25), ly + Inches(0.18),
                                Inches(0.4), Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

    inset_x = lx + Inches(0.25)
    add_text(slide, inset_x, ly + Inches(0.3), lw - Inches(0.5), Inches(0.3),
             f"MODEL · 0{role_num}", size=10, bold=True, color=ACCENT)
    add_text(slide, inset_x, ly + Inches(0.6), lw - Inches(0.5), Inches(0.5),
             name, size=22, bold=True, color=INK)

    cur_y = ly + Inches(1.25)

    def field(label, body, gap=0.55):
        nonlocal cur_y
        add_text(slide, inset_x, cur_y, Inches(2.0), Inches(0.3),
                 label.upper(), size=9, bold=True, color=INK_MUTED)
        add_text(slide, inset_x + Inches(2.0), cur_y, lw - Inches(2.5), Inches(0.5),
                 body, size=12, color=INK)
        cur_y += Inches(gap)

    field("Function", function)
    field("Base model", base)
    field("Dataset", dataset)
    field("Method", method, gap=0.7)
    field("Hardware", hardware)

    # Right card: metrics + hub
    rx = Inches(7.13); ry = Inches(2.3); rw = Inches(5.4); rh = Inches(4.5)
    add_rect(slide, rx, ry, rw, rh, fill=SURFACE_RAISED, line=INK_FAINT)

    # 4 stat boxes 2x2
    stats = [(time, "Wall-clock time"),
             (loss, "Final loss"),
             (cost, "Compute cost"),
             (size, "Adapter size")]
    sw_in = (rw - Inches(0.7)) / 2
    sh_in = Inches(1.0)
    for i, (val, label) in enumerate(stats):
        col = i % 2
        row = i // 2
        sx_box = rx + Inches(0.25) + (sw_in + Inches(0.2)) * col
        sy_box = ry + Inches(0.25) + (sh_in + Inches(0.15)) * row
        add_text(slide, sx_box, sy_box, sw_in, Inches(0.6),
                 val, size=28, bold=True, color=INK)
        add_text(slide, sx_box, sy_box + Inches(0.6), sw_in, Inches(0.4),
                 label.upper(), size=9, color=INK_MUTED)

    # Hub link
    add_text(slide, rx + Inches(0.25), ry + Inches(2.65), rw - Inches(0.5), Inches(0.3),
             "HUB", size=9, bold=True, color=INK_MUTED)
    add_text(slide, rx + Inches(0.25), ry + Inches(2.95), rw - Inches(0.5), Inches(0.5),
             hub, size=11, color=ACCENT, font=FONT_MONO)

    # Note
    add_text(slide, rx + Inches(0.25), ry + Inches(3.6), rw - Inches(0.5), Inches(0.8),
             note, size=11, color=INK_MUTED)

    add_footer(slide, f"Phase 4 · Model 0{role_num}", f"{page} / 13")


# Model 01 · SQL Generator
model_slide(
    num_label="08", page="08", role_num=1,
    name="SQL Generator",
    function="Translates a natural-language question and schema into a valid SQL query.",
    base="Qwen 2.5 Coder 7B Instruct",
    dataset="text-to-sql-mix-v2 (672,949 examples)",
    dataset_url="huggingface.co/datasets/DanielRegaladoCardoso/text-to-sql-mix-v2",
    method="QLoRA r=16, α=32, 4-bit base · TRL packing=True (154,462 packed sequences) · 1 epoch · 9,654 steps",
    hardware="1× NVIDIA L40S (48 GB) on Hugging Face Jobs",
    time="13.5h", loss="0.27", cost="~$24", size="161 MB",
    hub="DanielRegaladoCardoso/sql-generator-qwen25-coder-7b-lora",
    note="Sequence packing reduced training time from approximately 21 h to 13.5 h.",
)

# Model 02 · Chart Reasoner
model_slide(
    num_label="09", page="09", role_num=2,
    name="Chart Reasoner",
    function="Given a question and SQL result rows, decides the chart type and which columns to plot.",
    base="Microsoft Phi-3 Mini 4k Instruct",
    dataset="chart-reasoning-mix-v1 (~75,000 pairs)",
    dataset_url="huggingface.co/datasets/DanielRegaladoCardoso/chart-reasoning-mix-v1",
    method="QLoRA r=16, α=32, 4-bit base · 1 epoch · structured-JSON output objective",
    hardware="HF Jobs A10G / Colab Pro",
    time="~3h", loss="~0.31", cost="~$3", size="38 MB",
    hub="DanielRegaladoCardoso/chart-reasoner-phi3-mini-adapter-only",
    note="Outputs a JSON spec with chart_type, x_column, y_column, title, color, and rationale.",
)

# Model 03 · SVG Renderer
model_slide(
    num_label="10", page="10", role_num=3,
    name="SVG Renderer",
    function="Given a chart spec and data, produces inline SVG markup for the visualization.",
    base="DeepSeek Coder 1.3B Instruct",
    dataset="svg-chart-render-v1 (~25,000 chart-spec → SVG pairs)",
    dataset_url="huggingface.co/datasets/DanielRegaladoCardoso/svg-chart-render-v1",
    method="QLoRA r=16, α=32, 4-bit base · 1 epoch · code-generation objective",
    hardware="Colab T4",
    time="~2h", loss="~0.40", cost="~$1", size="22 MB",
    hub="DanielRegaladoCardoso/svg-renderer-deepseek-coder-1.3b-lora",
    note="When model output fails SVG validation, the system falls back to a themed Plotly renderer.",
)


# ============================================================ 11 ARCHITECTURE
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 5 — System architecture")
add_title(slide, "End-to-end query flow")

# Simple horizontal flow boxes
flow = [
    ("CSV / NL Q", False, False),
    ("Schema +\nDuckDB", False, False),
    ("Orchestrator", False, True),  # dark
    ("SQL Generator\nQwen + LoRA", True, False),
    ("Chart Reasoner\nPhi-3 + LoRA", True, False),
    ("SVG Renderer\nDeepSeek + LoRA", True, False),
    ("Chart +\nfinding", False, False, True),  # accent fill
]
fw = Inches(1.55); fg = Inches(0.18)
total = fw * len(flow) + fg * (len(flow) - 1)
fx = (SLIDE_W - total) / 2
for i, item in enumerate(flow):
    if len(item) == 4:
        label, accent_border, dark_fill, accent_fill = item
    else:
        label, accent_border, dark_fill = item
        accent_fill = False
    x = fx + (fw + fg) * i
    fill = INK if dark_fill else (ACCENT if accent_fill else SURFACE_RAISED)
    line = ACCENT if accent_border else INK_FAINT
    line_w = 1.5 if accent_border else 0.75
    text_color = WHITE if (dark_fill or accent_fill) else INK
    add_rect(slide, x, Inches(3.0), fw, Inches(1.5), fill=fill, line=line, line_w=line_w)
    add_text(slide, x, Inches(3.4), fw, Inches(0.9), label,
             size=11, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        ax = x + fw + Inches(0.04)
        add_line(slide, ax, Inches(3.75), ax + Inches(0.10), Inches(3.75),
                 color=INK_MUTED, weight=0.75)

# Description
add_text(slide, Inches(0.83), Inches(5.3), Inches(11.67), Inches(0.5),
         "Per query: four model invocations and a DuckDB execution. End-to-end latency: 5–8 seconds on a warm GPU.",
         size=14, color=INK)
add_text(slide, Inches(0.83), Inches(5.85), Inches(11.67), Inches(0.5),
         "All adapters loaded once at module level on a half-H200 via Hugging Face ZeroGPU.",
         size=13, color=INK_MUTED)
add_text(slide, Inches(0.83), Inches(6.3), Inches(11.67), Inches(0.5),
         "Self-correcting SQL: failed queries retried up to three times with the error in context.",
         size=13, color=INK_MUTED)

add_footer(slide, "Phase 5 · System architecture", "11 / 13")


# ============================================================ 12 COST
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Phase 6 — Deployment and cost")
add_title(slide, "Total compute cost: approximately $30")

# Left table
tx = Inches(0.83); ty = Inches(2.4); tw = Inches(6.5); th = Inches(4.0)
add_rect(slide, tx, ty, tw, th, fill=SURFACE_RAISED, line=INK_FAINT, line_w=0.75)

cost_rows = [
    ("Stage", "Compute", "Cost", True),
    ("SQL Generator training", "HF Jobs L40S, 13.5 h", "~$24", False),
    ("Chart Reasoner training", "Colab / HF Jobs", "~$3", False),
    ("SVG Renderer training", "Colab / HF Jobs", "~$1", False),
    ("GPT-4.1-nano distillation", "OpenAI Batch API", "~$2.50", False),
    ("Inference hosting", "HF Spaces ZeroGPU", "$0", False),
    ("Total", "", "~$30", "bold"),
]
for i, row in enumerate(cost_rows):
    if len(row) == 4:
        a, b, c, kind = row
    yy = ty + Inches(0.25 + i * 0.5)
    is_header = (kind is True)
    is_total = (kind == "bold")
    color = INK
    bold = is_header or is_total
    size_pt = 10 if is_header else 12
    add_text(slide, tx + Inches(0.25), yy, Inches(2.5), Inches(0.4),
             a, size=size_pt, bold=bold, color=INK_MUTED if is_header else color)
    add_text(slide, tx + Inches(2.85), yy, Inches(2.5), Inches(0.4),
             b, size=size_pt, bold=bold, color=INK_MUTED if is_header else color)
    add_text(slide, tx + Inches(5.4), yy, Inches(0.9), Inches(0.4),
             c, size=size_pt, bold=bold, color=INK_MUTED if is_header else color)
    if is_header or i == len(cost_rows) - 2:
        add_line(slide, tx + Inches(0.15), yy + Inches(0.42),
                 tx + tw - Inches(0.15), yy + Inches(0.42),
                 color=INK_FAINT, weight=0.5)

# Right
add_text(slide, Inches(7.8), Inches(2.6), Inches(4.7), Inches(1.5),
         "$30", size=84, bold=True, color=INK)
add_text(slide, Inches(7.8), Inches(4.0), Inches(4.7), Inches(0.4),
         "ALL-IN COMPUTE COST", size=10, bold=True, color=INK_MUTED)
add_text(slide, Inches(7.8), Inches(4.6), Inches(4.7), Inches(0.4),
         "WHY THE COST IS LOW", size=10, bold=True, color=ACCENT)
add_paragraphs(slide, Inches(7.8), Inches(4.95), Inches(4.7), Inches(2.0), [
    "•  Unsloth QLoRA reduces memory and time",
    "•  Sequence packing cuts SQL training time",
    "•  ZeroGPU is free for inference",
    "•  Only adapters are stored, not full models",
], size=11, color=INK)

add_footer(slide, "Phase 6 · Deployment and cost", "12 / 13")


# ============================================================ 13 DEMO + CONCLUSION
slide = prs.slides.add_slide(blank)
set_bg(slide, SURFACE)
add_accent_bar(slide)
add_kicker(slide, "Demonstration and conclusion")
add_title(slide, "Live system demo and summary")

# Demo strip
add_rect(slide, Inches(0.83), Inches(2.3), Inches(11.67), Inches(1.4),
         fill=SURFACE_RAISED, line=ACCENT, line_w=1.5)
add_text(slide, Inches(1.1), Inches(2.5), Inches(11.0), Inches(0.4),
         "LIVE DEMO", size=10, bold=True, color=ACCENT)
add_text(slide, Inches(1.1), Inches(2.85), Inches(11.0), Inches(0.5),
         "huggingface.co/spaces/DanielRegaladoCardoso/sql-agent",
         size=18, bold=True, color=INK)
add_text(slide, Inches(1.1), Inches(3.3), Inches(11.0), Inches(0.4),
         "Upload a CSV, ask a question. The system returns SQL, results, a chart, and a written finding.",
         size=12, color=INK_MUTED)

# Two columns
add_text(slide, Inches(0.83), Inches(4.0), Inches(5.5), Inches(0.4),
         "WHAT WE DELIVERED", size=10, bold=True, color=ACCENT)
add_paragraphs(slide, Inches(0.83), Inches(4.35), Inches(5.5), Inches(2.5), [
    "•  Three open-source datasets on Hugging Face",
    "•  Three QLoRA adapters trained on those datasets",
    "•  A working multi-model agent",
    "•  Total compute cost approximately $30",
], size=12, color=INK)

add_text(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.4),
         "FUTURE WORK", size=10, bold=True, color=ACCENT)
add_paragraphs(slide, Inches(7.0), Inches(4.35), Inches(5.5), Inches(2.5), [
    "•  Quantitative evaluation on Spider, WikiSQL, BIRD",
    "•  Multi-turn conversational memory",
    "•  Anomaly detection on uploaded data",
    "•  Statistical summary at ingestion",
], size=12, color=INK)

# Footer
add_text(slide, Inches(0.83), Inches(6.7), Inches(11.67), Inches(0.3),
         "github.com/DanielRegaladoUMiami/sql-agent-llmops",
         size=11, color=ACCENT, font=FONT_MONO)
add_text(slide, Inches(0.83), Inches(7.05), Inches(11.67), Inches(0.3),
         "Daniel Regalado Cardoso · Nefeli Zafeiri · Oliver Mazariegos · Eleniz Espina · MSBA UMiami 2026",
         size=10, color=INK_MUTED)


# ------------------------------------------------------------ save
out = "slides-editable.pptx"
prs.save(out)
print(f"Saved: {out}")
